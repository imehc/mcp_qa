"""
PowerPoint文档解析器
支持.pptx和.ppt格式的PowerPoint文件解析
"""

import os
import logging
from typing import Dict, Any, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain.text_splitter import RecursiveCharacterTextSplitter
import zipfile
import xml.etree.ElementTree as ET

from ..exceptions import ParsingError

logger = logging.getLogger(__name__)


class PowerPointParser:
    """PowerPoint文档解析器"""
    
    def __init__(self):
        """初始化PowerPoint解析器"""
        self.supported_extensions = ['.pptx', '.ppt']
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", "。", "！", "？", ";", ":", ".", "!", "?"]
        )
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否可以解析该文件"""
        return any(file_path.lower().endswith(ext) for ext in self.supported_extensions)
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        解析PowerPoint文件
        
        参数:
            file_path: PowerPoint文件路径
            **kwargs: 其他参数
                - extract_notes: 是否提取备注，默认True
                - extract_images: 是否提取图片信息，默认True
                - include_slide_numbers: 是否包含幻灯片编号，默认True
        
        返回:
            解析结果字典
        """
        if not os.path.exists(file_path):
            raise ParsingError(f"文件不存在: {file_path}")
        
        if not self.can_parse(file_path):
            raise ParsingError(f"不支持的文件格式: {file_path}")
        
        # 检查是否为.ppt格式（老版本）
        if file_path.lower().endswith('.ppt'):
            return self._parse_legacy_ppt(file_path, **kwargs)
        
        try:
            extract_notes = kwargs.get('extract_notes', True)
            extract_images = kwargs.get('extract_images', True)
            include_slide_numbers = kwargs.get('include_slide_numbers', True)
            
            logger.info(f"开始解析PowerPoint文件: {file_path}")
            
            # 加载演示文稿
            presentation = Presentation(file_path)
            
            slides_info = []
            all_text_content = []
            total_slides = len(presentation.slides)
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                
                # 提取幻灯片内容
                slide_content = self._extract_slide_content(
                    slide, 
                    slide_number, 
                    extract_notes=extract_notes,
                    extract_images=extract_images,
                    include_slide_numbers=include_slide_numbers
                )
                
                slides_info.append(slide_content)
                
                # 收集文本内容
                slide_text = self._format_slide_text(slide_content)
                all_text_content.append(slide_text)
            
            # 合并所有文本内容
            full_text = "\n\n".join(all_text_content)
            
            # 文本分块
            chunks = self.text_splitter.split_text(full_text) if full_text.strip() else []
            
            # 获取演示文稿元数据
            metadata = self._get_presentation_metadata(file_path)
            
            result = {
                "file_path": file_path,
                "file_type": "powerpoint",
                "total_slides": total_slides,
                "slides": slides_info,
                "full_text": full_text,
                "chunks": chunks,
                "chunk_count": len(chunks),
                "metadata": metadata
            }
            
            logger.info(f"PowerPoint文件解析完成: {total_slides}张幻灯片")
            return result
            
        except Exception as e:
            logger.error(f"PowerPoint文件解析失败: {file_path}, 错误: {e}")
            raise ParsingError(f"PowerPoint解析失败: {str(e)}")
    
    def _extract_slide_content(self, slide, slide_number: int, **kwargs) -> Dict[str, Any]:
        """提取单张幻灯片的内容"""
        extract_notes = kwargs.get('extract_notes', True)
        extract_images = kwargs.get('extract_images', True)
        include_slide_numbers = kwargs.get('include_slide_numbers', True)
        
        slide_content = {
            "slide_number": slide_number if include_slide_numbers else None,
            "title": "",
            "content": [],
            "notes": "",
            "images": [],
            "tables": [],
            "shapes": []
        }
        
        # 提取幻灯片标题和内容
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        # 判断是否为标题
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            try:
                                placeholder = shape.placeholder_format
                                if placeholder.type == 1:  # 标题占位符
                                    slide_content["title"] = text
                                else:
                                    slide_content["content"].append(text)
                            except:
                                slide_content["content"].append(text)
                        else:
                            slide_content["content"].append(text)
                
                # 提取图片信息
                if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_info = {
                            "type": "image",
                            "width": shape.width,
                            "height": shape.height,
                            "left": shape.left,
                            "top": shape.top
                        }
                        # 尝试获取图片名称
                        if hasattr(shape, 'image'):
                            image_info["filename"] = getattr(shape.image, 'filename', 'unknown')
                        slide_content["images"].append(image_info)
                    except Exception as e:
                        logger.debug(f"提取图片信息失败: {e}")
                
                # 提取表格
                if shape.has_table:
                    try:
                        table_data = self._extract_table_data(shape.table)
                        slide_content["tables"].append(table_data)
                    except Exception as e:
                        logger.debug(f"提取表格数据失败: {e}")
                
                # 记录其他形状类型
                shape_info = {
                    "type": str(shape.shape_type),
                    "has_text": shape.has_text_frame and bool(shape.text.strip())
                }
                slide_content["shapes"].append(shape_info)
                
            except Exception as e:
                logger.debug(f"处理形状时出错: {e}")
                continue
        
        # 提取备注
        if extract_notes and slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        notes_text = shape.text.strip()
                        if notes_text and not notes_text.startswith("单击此处"):
                            slide_content["notes"] = notes_text
                            break
            except Exception as e:
                logger.debug(f"提取备注失败: {e}")
        
        return slide_content
    
    def _extract_table_data(self, table) -> Dict[str, Any]:
        """提取表格数据"""
        table_data = {
            "rows": table.rows.__len__(),
            "columns": table.columns.__len__(),
            "data": []
        }
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_data.append(cell_text)
            table_data["data"].append(row_data)
        
        return table_data
    
    def _format_slide_text(self, slide_content: Dict[str, Any]) -> str:
        """将幻灯片内容格式化为文本"""
        lines = []
        
        if slide_content.get("slide_number"):
            lines.append(f"幻灯片 {slide_content['slide_number']}")
        
        if slide_content.get("title"):
            lines.append(f"标题: {slide_content['title']}")
        
        if slide_content.get("content"):
            lines.append("内容:")
            for content in slide_content["content"]:
                lines.append(f"- {content}")
        
        if slide_content.get("tables"):
            lines.append("表格:")
            for i, table in enumerate(slide_content["tables"]):
                lines.append(f"表格 {i+1}:")
                for row in table["data"]:
                    lines.append(" | ".join(row))
        
        if slide_content.get("notes"):
            lines.append(f"备注: {slide_content['notes']}")
        
        if slide_content.get("images"):
            lines.append(f"图片数量: {len(slide_content['images'])}")
        
        return "\n".join(lines)
    
    def _get_presentation_metadata(self, file_path: str) -> Dict[str, Any]:
        """获取演示文稿元数据"""
        metadata = {}
        
        try:
            # 从文件属性获取基本信息
            stat = os.stat(file_path)
            metadata.update({
                "file_size": stat.st_size,
                "created_time": stat.st_ctime,
                "modified_time": stat.st_mtime
            })
            
            # 尝试从PPTX文件中提取更多元数据
            if file_path.lower().endswith('.pptx'):
                try:
                    with zipfile.ZipFile(file_path, 'r') as zip_file:
                        # 读取核心属性
                        if 'docProps/core.xml' in zip_file.namelist():
                            core_xml = zip_file.read('docProps/core.xml')
                            root = ET.fromstring(core_xml)
                            
                            # 定义命名空间
                            ns = {
                                'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                                'dc': 'http://purl.org/dc/elements/1.1/',
                                'dcterms': 'http://purl.org/dc/terms/'
                            }
                            
                            # 提取元数据
                            title_elem = root.find('.//dc:title', ns)
                            if title_elem is not None:
                                metadata['title'] = title_elem.text
                            
                            creator_elem = root.find('.//dc:creator', ns)
                            if creator_elem is not None:
                                metadata['creator'] = creator_elem.text
                            
                            subject_elem = root.find('.//dc:subject', ns)
                            if subject_elem is not None:
                                metadata['subject'] = subject_elem.text
                            
                            description_elem = root.find('.//dc:description', ns)
                            if description_elem is not None:
                                metadata['description'] = description_elem.text
                
                except Exception as e:
                    logger.debug(f"提取PPTX元数据失败: {e}")
        
        except Exception as e:
            logger.debug(f"获取演示文稿元数据失败: {e}")
        
        return metadata
    
    def _parse_legacy_ppt(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """解析老版本.ppt文件"""
        try:
            # 对于.ppt文件，我们提供有限的支持
            # 可以尝试使用python-pptx的一些基本功能，但功能会受限
            
            logger.warning(f"解析.ppt文件支持有限: {file_path}")
            
            result = {
                "file_path": file_path,
                "file_type": "powerpoint_legacy",
                "total_slides": 0,
                "slides": [],
                "full_text": "",
                "chunks": [],
                "chunk_count": 0,
                "metadata": {
                    "note": "老版本.ppt文件支持有限，建议转换为.pptx格式"
                },
                "error": "Limited support for .ppt format",
                "suggestion": "请将.ppt文件转换为.pptx格式以获得完整的解析支持"
            }
            
            return result
            
        except Exception as e:
            logger.error(f".ppt文件解析失败: {e}")
            raise ParsingError(f".ppt文件解析失败: {str(e)}")
    
    def get_slide_count(self, file_path: str) -> int:
        """获取幻灯片数量"""
        try:
            if file_path.lower().endswith('.ppt'):
                return 0  # .ppt文件暂不支持
            
            presentation = Presentation(file_path)
            return len(presentation.slides)
            
        except Exception as e:
            logger.error(f"获取幻灯片数量失败: {e}")
            return 0
    
    def extract_slide_titles(self, file_path: str) -> List[str]:
        """提取所有幻灯片标题"""
        try:
            if file_path.lower().endswith('.ppt'):
                return []  # .ppt文件暂不支持
            
            presentation = Presentation(file_path)
            titles = []
            
            for slide in presentation.slides:
                title = ""
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        try:
                            if shape.placeholder_format.type == 1:  # 标题占位符
                                title = shape.text.strip()
                                break
                        except:
                            continue
                
                titles.append(title or f"幻灯片 {len(titles) + 1}")
            
            return titles
            
        except Exception as e:
            logger.error(f"提取幻灯片标题失败: {e}")
            return []


def create_powerpoint_parser() -> PowerPointParser:
    """创建PowerPoint解析器实例"""
    return PowerPointParser()


# 测试函数
def test_powerpoint_parser():
    """测试PowerPoint解析器"""
    parser = PowerPointParser()
    
    # 测试文件类型检查
    assert parser.can_parse("test.pptx")
    assert parser.can_parse("test.ppt")
    assert not parser.can_parse("test.pdf")
    
    logger.info("PowerPoint解析器测试通过")


if __name__ == "__main__":
    test_powerpoint_parser()